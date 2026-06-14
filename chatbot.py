import streamlit as st
import os, tempfile, uuid, asyncio, json, shutil, re
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from gtts import gTTS
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.http.models import Distance, VectorParams
from langchain_text_splitters import RecursiveCharacterTextSplitter
from fastembed import TextEmbedding

from agents import Agent, Runner

load_dotenv()

# ─── Paths ────────────────────────────────────────────────────────────────────
SPEECHES_DIR  = Path("speeches")          # all TTS audio files
HISTORY_FILE  = Path("chat_history.json") # persistent chat history
QDRANT_PATH   = "./qdrant_db"
SPEECHES_DIR.mkdir(exist_ok=True)

# ─── Helpers: persistent history ──────────────────────────────────────────────
def load_history() -> dict:
    if HISTORY_FILE.exists():
        try:
            return json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {}

def save_history(all_chats: dict):
    HISTORY_FILE.write_text(
        json.dumps(all_chats, ensure_ascii=True, indent=2),
        encoding="utf-8"
    )

# ─── Document loaders ─────────────────────────────────────────────────────────
def load_pdf(path: str) -> list[str]:
    """Extract text from PDF using pymupdf if available, else pdfminer."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        pages = []
        for page in doc:
            text = page.get_text()
            # Also try to extract tables as text
            if text.strip():
                pages.append(text)
        return pages
    except ImportError:
        pass
    try:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(path)
        docs = loader.load()
        return [d.page_content for d in docs]
    except Exception as e:
        return [f"[PDF load error: {e}]"]

def load_docx(path: str) -> list[str]:
    from docx import Document
    doc = Document(path)
    pages, chunk = [], []
    for para in doc.paragraphs:
        chunk.append(para.text)
        if len("\n".join(chunk)) > 1000:
            pages.append("\n".join(chunk))
            chunk = []
    if chunk:
        pages.append("\n".join(chunk))
    return pages

def safe_text(t: str) -> str:
    return t.encode("utf-8", errors="replace").decode("utf-8")

def load_pptx(path: str) -> list[str]:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(safe_text(shape.text.strip()))
        if texts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return slides

def load_txt(path: str) -> list[str]:
    content = Path(path).read_text(errors="replace")
    # split into ~1000 char chunks
    return [content[i:i+1000] for i in range(0, len(content), 900)]

def load_image_ocr(path: str) -> list[str]:
    from PIL import Image
    import pytesseract
    img = Image.open(path)
    text = pytesseract.image_to_string(img)
    return [text] if text.strip() else ["[No text found in image via OCR]"]

def detect_and_load(path: str, filename: str) -> list[str]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return load_pdf(path)
    elif ext in (".docx", ".doc"):
        return load_docx(path)
    elif ext in (".pptx", ".ppt"):
        return load_pptx(path)
    elif ext in (".txt", ".md", ".csv"):
        return load_txt(path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        return load_image_ocr(path)
    else:
        return load_txt(path)

# ─── Qdrant per-chat collection ───────────────────────────────────────────────
def get_collection_name(chat_id: str) -> str:
    return f"rag_{chat_id.replace('-', '')[:20]}"

def setup_qdrant(chat_id: str):
    client = QdrantClient(path=QDRANT_PATH)
    embedding_model = TextEmbedding()
    cname = get_collection_name(chat_id)
    try:
        client.create_collection(
            collection_name=cname,
            vectors_config=VectorParams(size=384, distance=Distance.COSINE)
        )
    except Exception:
        pass
    return client, embedding_model, cname

def delete_collection(chat_id: str):
    """Delete the vector store for a given chat (called on chat end)."""
    try:
        client = QdrantClient(path=QDRANT_PATH)
        cname = get_collection_name(chat_id)
        client.delete_collection(cname)
    except Exception:
        pass

def index_document(chat_id: str, pages: list[str]) -> int:
    client, model, cname = setup_qdrant(chat_id)
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=60)
    all_chunks = []
    for page in pages:
        if page.strip():
            all_chunks.extend(splitter.split_text(page))

    points = []
    for chunk in all_chunks:
        vector = list(model.embed([chunk]))[0].tolist()
        points.append(
            models.PointStruct(
                id=str(uuid.uuid4()),
                vector=vector,
                payload={"content": chunk}
            )
        )
    if points:
        client.upsert(collection_name=cname, points=points)
    return len(all_chunks)

def search_qdrant(chat_id: str, query: str, limit: int = 5) -> str:
    try:
        client, model, cname = setup_qdrant(chat_id)
        query_vector = list(model.embed([query]))[0].tolist()
        results = client.query_points(
            collection_name=cname,
            query=query_vector,
            limit=limit
        ).points
        return "\n\n---\n".join([r.payload['content'] for r in results if r.payload])
    except Exception as e:
        return f"[Vector search error: {e}]"

# ─── Session state init ───────────────────────────────────────────────────────
def init_session():
    if "all_chats" not in st.session_state:
        st.session_state.all_chats = load_history()

    if "current_chat_id" not in st.session_state:
        if st.session_state.all_chats:
            # resume the most recent chat
            st.session_state.current_chat_id = list(st.session_state.all_chats.keys())[-1]
        else:
            new_id = str(uuid.uuid4())
            st.session_state.current_chat_id = new_id
            st.session_state.all_chats[new_id] = {
                "name": "New Chat",
                "messages": [],
                "doc_name": None,
                "doc_indexed": False,
                "ended": False,
                "created_at": datetime.now().isoformat()
            }
            save_history(st.session_state.all_chats)

def new_chat():
    new_id = str(uuid.uuid4())
    st.session_state.all_chats[new_id] = {
        "name": "New Chat",
        "messages": [],
        "doc_name": None,
        "doc_indexed": False,
        "ended": False,
        "created_at": datetime.now().isoformat()
    }
    st.session_state.current_chat_id = new_id
    save_history(st.session_state.all_chats)
    st.rerun()

def end_chat(chat_id: str):
    """Mark chat as ended and delete vector DB for it."""
    if chat_id in st.session_state.all_chats:
        st.session_state.all_chats[chat_id]["ended"] = True
        save_history(st.session_state.all_chats)
        delete_collection(chat_id)
        st.success("✅ Chat ended. Document memory cleared. History saved.")

# ─── TTS ──────────────────────────────────────────────────────────────────────
def speak(text: str) -> str:
    clean = re.sub(r"[#*`_>]", "", text)[:800]
    audio_path = SPEECHES_DIR / f"speech_{uuid.uuid4()}.mp3"
    tts = gTTS(text=clean, lang='en')
    tts.save(str(audio_path))
    return str(audio_path)

# ─── AI Call ──────────────────────────────────────────────────────────────────
def run_agent(system: str, prompt: str) -> str:
    agent = Agent(name="TeacherAI", instructions=system)
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    result = loop.run_until_complete(Runner.run(agent, prompt))
    return result.final_output

# ─── Build system prompt ──────────────────────────────────────────────────────
TEACHER_SYSTEM = """You are an expert AI teacher. Your ONLY knowledge source is the document context provided below.

STRICT RULES:
1. Answer ONLY from the provided document context. Do NOT use outside knowledge.
2. If the answer is not in the context, say: "This information is not in the uploaded document."
3. For numerical problems/formulas: show step-by-step working clearly.
4. For graphs/images described in text: explain them clearly with words.
5. Be a patient teacher — explain concepts clearly as if teaching a student.
6. For exam questions: give important questions with model answers from the document.
7. Always cite which part of the document your answer comes from.
8. For "teach me this topic", give a structured lesson plan with key points from the doc.
"""

def build_prompt(context: str, history: list, query: str) -> str:
    mem = "\n".join([f"{m['role'].upper()}: {m['content']}" for m in history[-6:]])
    return f"""DOCUMENT CONTEXT:
{context}

CONVERSATION HISTORY:
{mem}

STUDENT QUESTION: {query}

Answer strictly based on the document context above:"""

# ─── Main App ─────────────────────────────────────────────────────────────────
def main():
    st.set_page_config(page_title="AI Doc Teacher", page_icon="📚", layout="wide")
    init_session()

    # ── SIDEBAR ──────────────────────────────────────────────────────────────
    with st.sidebar:
        st.title("📚 AI Doc Teacher")
        st.caption("Upload any doc — I'll teach it to you!")

        if st.button("➕ New Chat", use_container_width=True, type="primary"):
            new_chat()

        st.divider()
        st.subheader("💬 Chat History")

        for chat_id, data in reversed(list(st.session_state.all_chats.items())):
            is_active = (chat_id == st.session_state.current_chat_id)
            ended_tag = " 🔒" if data.get("ended") else ""
            doc_tag   = f" 📄" if data.get("doc_name") else ""
            label     = f"{data['name']}{doc_tag}{ended_tag}"
            col1, col2 = st.columns([5, 1])
            with col1:
                if st.button(label, key=f"nav_{chat_id}", use_container_width=True,
                             type="primary" if is_active else "secondary"):
                    st.session_state.current_chat_id = chat_id
                    st.rerun()
            with col2:
                if st.button("🗑", key=f"del_{chat_id}", help="Delete chat"):
                    delete_collection(chat_id)
                    del st.session_state.all_chats[chat_id]
                    save_history(st.session_state.all_chats)
                    if is_active:
                        if st.session_state.all_chats:
                            st.session_state.current_chat_id = list(st.session_state.all_chats.keys())[-1]
                        else:
                            new_chat()
                    st.rerun()

        st.divider()
        st.subheader("🎙️ Saved Speeches")

        speech_files = sorted(SPEECHES_DIR.glob("*.mp3"), key=lambda f: f.stat().st_mtime, reverse=True)

        if not speech_files:
            st.caption("No speeches yet. Ask a question to generate one!")
        else:
            st.caption(f"{len(speech_files)} recording(s) stored")
            for sp in speech_files:
                mtime = datetime.fromtimestamp(sp.stat().st_mtime).strftime("%b %d, %H:%M")
                with st.expander(f"🔊 {mtime}", expanded=False):
                    st.audio(str(sp))
                    if st.button("🗑 Delete", key=f"del_sp_{sp.name}"):
                        sp.unlink()
                        st.rerun()

            if st.button("🗑 Clear All Speeches", use_container_width=True):
                for sp in speech_files:
                    sp.unlink()
                st.success("All speeches cleared!")
                st.rerun()

    # ── MAIN ─────────────────────────────────────────────────────────────────
    current_id   = st.session_state.current_chat_id
    current_chat = st.session_state.all_chats[current_id]

    st.title("📖 AI Document Teacher")

    # Ended chat banner
    if current_chat.get("ended"):
        st.warning("🔒 This chat has ended. Document memory cleared. You can still read history or start a **New Chat**.")
        for msg in current_chat["messages"]:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])
        return

    # ── Document Upload ───────────────────────────────────────────────────────
    st.subheader("📂 Upload Your Document")
    accepted = ["pdf", "docx", "doc", "pptx", "ppt", "txt", "md", "png", "jpg", "jpeg"]
    uploaded = st.file_uploader(
        "Upload PDF, Word, PPT, Notes, Image (OCR), or Text file",
        type=accepted,
        key=f"upload_{current_id}"
    )

    if uploaded and not current_chat["doc_indexed"]:
        with st.spinner(f"📥 Reading & indexing `{uploaded.name}`..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(uploaded.name).suffix) as tf:
                tf.write(uploaded.getbuffer())
                tmp_path = tf.name

            pages = detect_and_load(tmp_path, uploaded.name)
            os.unlink(tmp_path)

            n_chunks = index_document(current_id, pages)

            current_chat["doc_name"]    = uploaded.name
            current_chat["doc_indexed"] = True
            current_chat["name"]        = f"📄 {uploaded.name[:20]}"
            save_history(st.session_state.all_chats)

        st.success(f"✅ **{uploaded.name}** indexed! ({n_chunks} chunks) Now ask me anything about it.")

        # Ask user what they want to do with this doc
        st.info("💡 **What would you like to do?**\n\n"
                "- *Teach me this document from scratch*\n"
                "- *Give me the 10 most important exam questions*\n"
                "- *Explain all formulas/numericals*\n"
                "- *Summarize the key topics*\n"
                "- *Quiz me on this document*")

    if current_chat.get("doc_name") and current_chat["doc_indexed"]:
        st.caption(f"📄 Active document: **{current_chat['doc_name']}**")

    # ── Chat History display ──────────────────────────────────────────────────
    for msg in current_chat["messages"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # ── End chat button ───────────────────────────────────────────────────────
    col_a, col_b = st.columns([6, 1])
    with col_b:
        if st.button("🔚 End Chat", help="End session & clear document memory"):
            end_chat(current_id)
            st.rerun()

    # ── Chat Input ────────────────────────────────────────────────────────────
    query = st.chat_input("Ask about your document... (e.g. 'Teach me Chapter 2' / 'Give 10 exam questions')")

    if query:
        # Guard: no doc uploaded
        if not current_chat.get("doc_indexed"):
            st.warning("⚠️ Please upload a document first before asking questions.")
            return

        # Name the chat
        if current_chat["name"].startswith("📄") or current_chat["name"] == "New Chat":
            pass
        if len(current_chat["messages"]) == 0:
            short = query[:28] + ("..." if len(query) > 28 else "")
            current_chat["name"] = f"📄 {short}"

        with st.chat_message("user"):
            st.markdown(query)
        current_chat["messages"].append({"role": "user", "content": query})

        with st.spinner("🔍 Searching document & thinking..."):
            # Retrieve relevant context — more chunks for exam/teach requests
            limit = 8 if any(k in query.lower() for k in ["exam", "questions", "teach", "all", "formula", "summary"]) else 5
            context = search_qdrant(current_id, query, limit=limit)

            prompt = build_prompt(context, current_chat["messages"], query)

            # Detect exam-mode
            if re.search(r"(exam|important questions|question.*exam|quiz me|test me)", query, re.I):
                system = TEACHER_SYSTEM + "\n\nSPECIAL: User wants exam preparation. List exactly 10 most important questions with model answers from the document. Format as numbered list."
            elif re.search(r"(teach me|explain|lesson|concept|what is|how does)", query, re.I):
                system = TEACHER_SYSTEM + "\n\nSPECIAL: Teach this topic step-by-step with clear headings. Use examples from the document."
            elif re.search(r"(formula|numerical|calculate|solve|derivation|equation)", query, re.I):
                system = TEACHER_SYSTEM + "\n\nSPECIAL: Show all formulas and step-by-step numerical solutions clearly. Use proper notation."
            elif re.search(r"(graph|figure|image|diagram|table|chart)", query, re.I):
                system = TEACHER_SYSTEM + "\n\nSPECIAL: Describe any graphs, figures or tables from the document context in detail."
            else:
                system = TEACHER_SYSTEM

            response = run_agent(system, prompt)

        current_chat["messages"].append({"role": "assistant", "content": response})
        save_history(st.session_state.all_chats)

        with st.chat_message("assistant"):
            st.markdown(response)

        # TTS
        try:
            audio_path = speak(response)
            st.audio(audio_path, autoplay=True)
        except Exception as e:
            st.caption(f"🔇 TTS unavailable: {e}")

        st.rerun()

if __name__ == "__main__":
    main()